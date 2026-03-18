"""
Flask チャットボットアプリ
ChatGPT API を使用した対話アプリ
"""
import os
import sys
import json
import time
import re
import urllib.request
import urllib.error
import urllib.parse

# Windows で日本語を扱うときの ASCII エンコードエラーを防ぐ
if sys.platform == "win32":
    import io
    for name in ("stdout", "stderr"):
        stream = getattr(sys, name)
        if hasattr(stream, "buffer"):
            setattr(sys, name, io.TextIOWrapper(stream.buffer, encoding="utf-8", errors="replace"))

from flask import Flask, render_template, request, jsonify

# PDF / DOCX / PPTX 用（オプション：ライブラリが無い場合は該当形式をスキップ）
try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

app = Flask(__name__)
# 日本語などをそのまま JSON で返すため（ASCII に変換しない）
app.config["JSON_AS_ASCII"] = False

# コンテキスト用フォルダ（app.py と同じ場所の context/）
CONTEXT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "context")
CONTEXT_MAX_CHARS = 30000  # トークン制限を考慮した上限
CONTEXT_EXTENSIONS = (".txt", ".md", ".pdf", ".docx", ".pptx")

# 常に含める固定コンテキスト（Xアカウント情報など）
FIXED_CONTEXT = "X（旧Twitter）の @threee_sales はスリーグッドの田中祐貴のアカウントである。"

# 用途別プロンプト（Gem 風）の設定ファイル
PROMPTS_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "prompts.json")


def load_presets():
    """prompts.json から用途別プロンプトを読み込む。失敗時は標準のみ"""
    default_presets = [{"id": "default", "name": "標準", "prompt": ""}]
    if not os.path.isfile(PROMPTS_FILE):
        return default_presets
    try:
        with open(PROMPTS_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
        presets = data.get("presets") or default_presets
        return [p for p in presets if p.get("id") and p.get("name") is not None]
    except Exception:
        return default_presets


def get_preset_prompt(preset_id):
    """preset_id に対応するプロンプト文を返す。なければ空文字"""
    if not preset_id or not (preset_id := str(preset_id).strip()):
        return ""
    for p in load_presets():
        if (p.get("id") or "").strip() == preset_id:
            return (p.get("prompt") or "").strip()
    return ""


def save_preset_prompt(preset_id, prompt):
    """prompts.json の指定プリセットの prompt を更新する。成功で True"""
    preset_id = (preset_id or "").strip()
    if not preset_id or preset_id == "default":
        return False
    prompt = (prompt or "").strip() if prompt is not None else ""
    try:
        if not os.path.isfile(PROMPTS_FILE):
            return False
        with open(PROMPTS_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
        presets = data.get("presets") or []
        for p in presets:
            if (p.get("id") or "").strip() == preset_id:
                p["prompt"] = prompt
                break
        else:
            return False
        with open(PROMPTS_FILE, "w", encoding="utf-8") as f:
            json.dump({"presets": presets}, f, ensure_ascii=False, indent=2)
        return True
    except Exception:
        return False


def _extract_text_from_pdf(path):
    """PDF からテキストを抽出する"""
    if PdfReader is None:
        return None
    reader = PdfReader(path)
    parts = []
    for page in reader.pages:
        t = page.extract_text()
        if t:
            parts.append(t)
    return "\n".join(parts) if parts else ""


def _extract_text_from_docx(path):
    """DOCX からテキストを抽出する"""
    if DocxDocument is None:
        return None
    doc = DocxDocument(path)
    parts = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(parts) if parts else ""


def _extract_text_from_pptx(path):
    """PPTX からテキストを抽出する"""
    if Presentation is None:
        return None
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text)
    return "\n".join(parts) if parts else ""


def _read_file_text(path):
    """拡張子に応じてファイルからテキストを読み込む。失敗時は None"""
    name = os.path.basename(path)
    lower = name.lower()
    try:
        if lower.endswith(".txt") or lower.endswith(".md"):
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
        if lower.endswith(".pdf"):
            return _extract_text_from_pdf(path)
        if lower.endswith(".docx"):
            return _extract_text_from_docx(path)
        if lower.endswith(".pptx"):
            return _extract_text_from_pptx(path)
    except Exception:
        pass
    return None


def get_context_text():
    """context フォルダ内の .txt / .md / .pdf / .docx / .pptx を更新日時の新しい順に読み込み、1つの文字列にする"""
    if not os.path.isdir(CONTEXT_DIR):
        return ""
    parts = []
    try:
        files = []
        for name in os.listdir(CONTEXT_DIR):
            if any(name.lower().endswith(ext) for ext in CONTEXT_EXTENSIONS):
                path = os.path.join(CONTEXT_DIR, name)
                if os.path.isfile(path):
                    files.append((path, os.path.getmtime(path)))
        files.sort(key=lambda x: -x[1])  # 新しい順
        total = 0
        for path, _ in files:
            if total >= CONTEXT_MAX_CHARS:
                break
            text = _read_file_text(path)
            if text is None or not text.strip():
                continue
            name = os.path.basename(path)
            chunk = f"\n--- {name} ---\n{text}\n"
            if total + len(chunk) > CONTEXT_MAX_CHARS:
                chunk = chunk[: CONTEXT_MAX_CHARS - total]
            parts.append(chunk)
            total += len(chunk)
    except Exception:
        pass
    return "".join(parts)


# 環境変数 OPENAI_API_KEY からAPIキーを取得
def get_api_key():
    api_key = os.environ.get("OPENAI_API_KEY")
    if not api_key:
        raise ValueError("OPENAI_API_KEY が設定されていません。")
    return api_key


# 環境変数 GEMINI_API_KEY（Google AI Studio の API キー）
def get_gemini_api_key():
    api_key = os.environ.get("GEMINI_API_KEY")
    if not api_key:
        raise ValueError("GEMINI_API_KEY が設定されていません。")
    return api_key


def is_gemini_model(model):
    """選択されたモデルが Google Gemini かどうか"""
    if not model or not isinstance(model, str):
        return False
    return model.strip().lower().startswith("gemini-")


# 使用するLLMモデル（環境変数 OPENAI_CHAT_MODEL で変更可能。未設定時は gpt-4o-mini）
def get_chat_model():
    return os.environ.get("OPENAI_CHAT_MODEL", "gpt-4o-mini").strip() or "gpt-4o-mini"


def call_gemini_api(messages, api_key, model):
    """Google Gemini API を呼び出す。messages は OpenAI 形式 [{"role":"system|user|assistant","content":"..."}]"""
    model = (model or "gemini-2.0-flash").strip()
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key}"
    system_parts = []
    contents = []
    for m in messages:
        role = (m.get("role") or "").strip().lower()
        content = (m.get("content") or "").strip()
        if not content:
            continue
        if role == "system":
            system_parts.append(content)
            continue
        gemini_role = "model" if role == "assistant" else "user"
        contents.append({"role": gemini_role, "parts": [{"text": content}]})
    body = {
        "contents": contents,
        "generationConfig": {"temperature": 0.7},
    }
    if system_parts:
        body["systemInstruction"] = {"parts": [{"text": "\n\n".join(system_parts)}]}
    body_bytes = json.dumps(body, ensure_ascii=False).encode("utf-8")
    req = urllib.request.Request(
        url,
        data=body_bytes,
        method="POST",
        headers={"Content-Type": "application/json; charset=utf-8"},
    )
    with urllib.request.urlopen(req) as res:
        data = json.loads(res.read().decode("utf-8"))
    if "candidates" not in data or not data["candidates"]:
        raise RuntimeError(data.get("error", {}).get("message", "Gemini が応答を返しませんでした。") or str(data))
    parts = data["candidates"][0].get("content", {}).get("parts", [])
    if not parts:
        text = ""
    else:
        text = (parts[0].get("text") or "").strip()
    usage = _gemini_usage_from_response(data)
    return text, usage


def _gemini_usage_from_response(data):
    """Gemini API レスポンスから利用量を抽出。input/output/total トークン数"""
    um = data.get("usageMetadata") or data.get("usage_metadata") or {}
    prompt = um.get("promptTokenCount") or um.get("prompt_token_count") or 0
    candidates = um.get("candidatesTokenCount") or um.get("candidates_token_count") or 0
    total = um.get("totalTokenCount") or um.get("total_token_count") or (prompt + candidates)
    return {"input_tokens": prompt, "output_tokens": candidates, "total_tokens": total}


def call_chatgpt_api(messages, api_key, model=None):
    """UTF-8 で明示的にリクエストを送り、Windows の ASCII エンコードエラーを防ぐ"""
    url = "https://api.openai.com/v1/chat/completions"
    body = {
        "model": (model or get_chat_model()).strip() or get_chat_model(),
        "messages": messages,
        "temperature": 0.7,
    }
    # 日本語をそのまま送るため ensure_ascii=False、バイト列は UTF-8 で作成
    body_bytes = json.dumps(body, ensure_ascii=False).encode("utf-8")
    req = urllib.request.Request(
        url,
        data=body_bytes,
        method="POST",
        headers={
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json; charset=utf-8",
        },
    )
    with urllib.request.urlopen(req) as res:
        data = json.loads(res.read().decode("utf-8"))
    content = data["choices"][0]["message"]["content"]
    usage = _openai_usage_from_response(data)
    return content, usage


def _openai_usage_from_response(data):
    """OpenAI API レスポンスから利用量を抽出"""
    u = data.get("usage") or {}
    prompt = u.get("prompt_tokens", 0)
    completion = u.get("completion_tokens", 0)
    total = u.get("total_tokens", 0) or (prompt + completion)
    return {"input_tokens": prompt, "output_tokens": completion, "total_tokens": total}


@app.route("/")
def index():
    """チャット画面を表示"""
    return render_template("index.html")


@app.route("/api/context", methods=["GET"])
def api_context():
    """現在読み込まれているコンテキスト（ファイル一覧と先頭のプレビュー）を返す"""
    text = get_context_text()
    files = []
    if os.path.isdir(CONTEXT_DIR):
        try:
            for name in sorted(os.listdir(CONTEXT_DIR)):
                if any(name.lower().endswith(ext) for ext in CONTEXT_EXTENSIONS):
                    path = os.path.join(CONTEXT_DIR, name)
                    if os.path.isfile(path):
                        files.append({"name": name, "mtime": os.path.getmtime(path)})
            files.sort(key=lambda x: -x["mtime"])
        except Exception:
            pass
    preview = text[:500] + "…" if len(text) > 500 else text
    return jsonify({
        "files": [f["name"] for f in files],
        "preview": preview,
        "length": len(text),
        "model": get_chat_model(),
    })


@app.route("/api/presets", methods=["GET"])
def api_presets():
    """用途別プロンプト（Gem 風）の一覧を返す"""
    presets = [{"id": p.get("id", ""), "name": p.get("name", "")} for p in load_presets()]
    return jsonify({"presets": presets})


@app.route("/api/preset/<preset_id>", methods=["GET"])
def api_preset_detail(preset_id):
    """指定した用途のプロンプト全文を返す（選択時に読み込んで表示用）"""
    preset_id = (preset_id or "").strip()
    for p in load_presets():
        if (p.get("id") or "").strip() == preset_id:
            return jsonify({
                "id": p.get("id", ""),
                "name": p.get("name", ""),
                "prompt": (p.get("prompt") or "").strip(),
            })
    return jsonify({"id": preset_id, "name": "", "prompt": ""})


@app.route("/api/preset/<preset_id>", methods=["PUT"])
def api_preset_update(preset_id):
    """指定した用途のプロンプトを prompts.json に保存する"""
    preset_id = (preset_id or "").strip()
    if preset_id == "default":
        return jsonify({"error": "標準プリセットは編集できません"}), 400
    data = request.get_json() or {}
    prompt = (data.get("prompt") or "").strip() if data.get("prompt") is not None else ""
    if not save_preset_prompt(preset_id, prompt):
        return jsonify({"error": "保存に失敗しました"}), 500
    return jsonify({"ok": True, "id": preset_id})


@app.route("/api/chat", methods=["POST"])
def chat():
    """ユーザーメッセージを受け取り、ChatGPTの応答を返す"""
    data = request.get_json()
    if not data or "message" not in data:
        return jsonify({"error": "message が必要です"}), 400

    user_message = data["message"]
    history = data.get("history", [])  # 会話履歴（オプション）
    extra_context = (data.get("extra_context") or "").strip()  # 画面から渡す追加コンテキスト
    model_override = (data.get("model") or "").strip()  # 画面で選択したモデル（任意）
    preset_id = (data.get("preset_id") or "").strip()  # 用途別プロンプト（Gem 風）
    prompt_override = (data.get("prompt_override") or "").strip()  # UIで編集したプロンプト（あれば優先）

    use_gemini = is_gemini_model(model_override)
    try:
        api_key = get_gemini_api_key() if use_gemini else get_api_key()
    except ValueError as e:
        return jsonify({"error": str(e)}), 500

    # コンテキストをシステムメッセージとして先頭に付与
    file_context = get_context_text()
    context_parts = ["【常に参照する情報】\n" + FIXED_CONTEXT]
    preset_prompt = prompt_override if prompt_override else get_preset_prompt(preset_id)
    if preset_prompt:
        context_parts.append("【用途・指示】\n" + preset_prompt)
    if file_context:
        context_parts.append("【参考情報（フォルダから読み込み）】\n" + file_context)
    if extra_context:
        context_parts.append("【この会話で追加された参考情報】\n" + extra_context)
    if context_parts:
        system_content = (
            "【重要】以下にコンテキスト情報を記載します。\n"
            "・質問がコンテキストに記載されている内容（会社概要・議事録・会話メモ・Xのやり取り等）に関係する場合は、必ずコンテキストを読み込み、その内容を参照して回答すること。\n"
            "・コンテキストにない話題や一般論の質問の場合はその限りではない。\n"
            "・「用途・指示」がある場合は、その役割・トーンに従って応答すること。\n\n"
            "--- コンテキスト ---\n\n"
            + "\n\n".join(context_parts)
        )
    else:
        system_content = "ユーザーの質問に丁寧に答えてください。"

    messages = [{"role": "system", "content": system_content}]
    for h in history:
        messages.append({"role": "user", "content": h.get("user", "")})
        messages.append({"role": "assistant", "content": h.get("assistant", "")})
    messages.append({"role": "user", "content": user_message})

    try:
        if use_gemini:
            assistant_message, usage = call_gemini_api(messages, api_key, model_override)
        else:
            assistant_message, usage = call_chatgpt_api(messages, api_key, model=model_override or None)
        return jsonify({"reply": assistant_message, "usage": usage})
    except urllib.error.HTTPError as e:
        err_body = e.read().decode("utf-8", errors="replace")
        return jsonify({"error": f"APIエラー: {err_body}"}), 500
    except Exception as e:
        return jsonify({"error": f"APIエラー: {str(e)}"}), 500


# ---- Slack RAG ----

@app.route("/slack_rag")
def slack_rag_page():
    """Slack RAG の操作画面"""
    return render_template("slack_rag.html")




@app.route("/slack_rag_query", methods=["POST"])
def slack_rag_query():
    """
    Slack RAG に問い合わせる。
    リクエスト JSON: {"question": "...", "top_k": 8}
    レスポンス JSON: {"answer": "...", "sources": [...]}
    """
    try:
        from slack_rag.query_rag import answer as rag_answer
    except ImportError as e:
        return jsonify({"error": f"slack_rag モジュールが見つかりません: {e}"}), 500

    data = request.get_json() or {}
    question = (data.get("question") or "").strip()
    if not question:
        return jsonify({"error": "question が必要です"}), 400
    top_k = data.get("top_k")

    try:
        result = rag_answer(question, top_k=top_k)
        return jsonify(result)
    except ValueError as e:
        return jsonify({"error": str(e)}), 400
    except Exception as e:
        return jsonify({"error": f"RAG エラー: {str(e)}"}), 500


@app.route("/slack_rag_fetch", methods=["POST"])
def slack_rag_fetch():
    """
    Slack メッセージを取得して DB に保存する。
    リクエスト JSON: {"limit": 200}
    """
    try:
        from slack_rag.fetch_slack_messages import fetch_and_store, count_messages
    except ImportError as e:
        return jsonify({"error": f"slack_rag モジュールが見つかりません: {e}"}), 500

    data = request.get_json() or {}
    limit = data.get("limit")

    try:
        saved = fetch_and_store(limit=limit)
        total = count_messages()
        return jsonify({"saved": saved, "total": total})
    except ValueError as e:
        return jsonify({"error": str(e)}), 400
    except Exception as e:
        return jsonify({"error": f"Slack 取得エラー: {str(e)}"}), 500


@app.route("/slack_rag_build", methods=["POST"])
def slack_rag_build():
    """
    ChromaDB ベクトルストアを構築・更新する。
    リクエスト JSON: {"reset": false}
    """
    try:
        from slack_rag.build_vector_store import build
    except ImportError as e:
        return jsonify({"error": f"slack_rag モジュールが見つかりません: {e}"}), 500

    data = request.get_json() or {}
    reset = bool(data.get("reset", False))

    try:
        upserted = build(reset=reset)
        return jsonify({"upserted": upserted})
    except ValueError as e:
        return jsonify({"error": str(e)}), 400
    except Exception as e:
        return jsonify({"error": f"ベクトルストア構築エラー: {str(e)}"}), 500


if __name__ == "__main__":
    host = os.environ.get("FLASK_HOST", "0.0.0.0")
    # Railway は PORT 環境変数を自動設定する。FLASK_PORT でも上書き可能
    port = int(os.environ.get("PORT", os.environ.get("FLASK_PORT", "5000")))
    debug = os.environ.get("FLASK_DEBUG", "false").lower() == "true"
    app.run(host=host, debug=debug, port=port)
